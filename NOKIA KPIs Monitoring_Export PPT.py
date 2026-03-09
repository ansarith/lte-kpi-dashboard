import os, sys
import streamlit as st
import pandas as pd
import plotly.graph_objects as go
import plotly.express as px
import io
from pptx import Presentation
from pptx.util import Inches

st.set_page_config(layout="wide")

# ----------------------------------------------------
st.write("🚀 Running file:", os.path.abspath(__file__))
st.write("🟢 Python executable:", sys.executable)
# ----------------------------------------------------

# --- Load data ---
@st.cache_data
def load_data(path):
    df = pd.read_excel(path)
    df["Period start time"] = pd.to_datetime(df["Period start time"], errors="coerce")

    percentage_kpis = [col for col in df.columns if "%" in col or "Rate" in col]

    for col in percentage_kpis:
        if col in df.columns and pd.api.types.is_numeric_dtype(df[col]):
            if df[col].max() <= 1.0:
                df[col] = df[col] * 100

    return df


DATA_PATH = r"D:\Cell Database RAN Team 2024\NOKIA KPIs Coding Monitoring\4G_Main_KPIs_Report_SRAN21B-Sarith-2025_10_13-Site KCH2567RBR & 2070_LB.xlsx"
df = load_data(DATA_PATH)

st.title("📊 LTE KPI Dashboard")

# --- KPI selection ---
kpi_columns = [col for col in df.columns if col not in ["Period start time","LNBTS name","LNCEL name"]]

selected_kpis = st.multiselect(
    "Select KPI(s)",
    options=kpi_columns,
    default=kpi_columns[:4]
)

# --- Site filter ---
enodeb_selected = st.multiselect(
    "Select LNBTS name",
    options=sorted(df["LNBTS name"].unique())
)

# --- Cell filter ---
if enodeb_selected:
    cell_options = sorted(df[df["LNBTS name"].isin(enodeb_selected)]["LNCEL name"].unique())
else:
    cell_options = sorted(df["LNCEL name"].unique())

cell_selected = st.multiselect(
    "Select LNCEL name",
    options=cell_options
)

# --- Options ---
daily_option = st.checkbox("📅 Daily Aggregation")
group_option = st.checkbox("🏙️ Group by Site")

# --- Filter dataframe ---
plot_df = df.copy()

if enodeb_selected:
    plot_df = plot_df[plot_df["LNBTS name"].isin(enodeb_selected)]

if cell_selected:
    plot_df = plot_df[plot_df["LNCEL name"].isin(cell_selected)]


# --- Aggregation ---
def aggregate_data(df, kpis, daily=False, group=False):

    for kpi in kpis:
        df[kpi] = pd.to_numeric(df[kpi], errors="coerce")

    agg_dict = {}

    for kpi in kpis:

        if kpi in [
            "PDCP SDU Volume, DL",
            "PDCP SDU Volume, UL",
            "Total LTE data volume, DL + UL",
            "Avg RRC conn UE",
            "RRC Connected UEs Max (M8051C56)"
        ]:
            agg_dict[kpi] = "sum"

        else:
            agg_dict[kpi] = "mean"

    if daily:
        df["Date"] = df["Period start time"].dt.normalize()
        time_col = "Date"
    else:
        time_col = "Period start time"

    if not group:

        group_cols = [time_col]

        if "LNCEL name" in df.columns:
            group_cols.append("LNCEL name")

        grouped = df.groupby(group_cols, as_index=False).agg(agg_dict)

    else:

        grouped = df.groupby([time_col], as_index=False).agg(agg_dict)

    return grouped


plot_df = aggregate_data(plot_df, selected_kpis, daily_option, group_option)
# After aggregation
time_col = "Date" if daily_option else "Period start time"
plot_df[time_col] = pd.to_datetime(plot_df[time_col], errors="coerce")
plot_df = plot_df.dropna(subset=[time_col])

# Convert to formatted string for export/plot
if daily_option:
    plot_df["Time_str"] = plot_df[time_col].dt.strftime("%Y-%m-%d")
else:
    plot_df["Time_str"] = plot_df[time_col].dt.strftime("%Y-%m-%d %H:%M")

# --- Dashboard ---
if not plot_df.empty:

    colors = px.colors.qualitative.Dark24
    figures = []

    cols = st.columns(2)

    for idx, selected_kpi in enumerate(selected_kpis[:4]):

        fig = go.Figure()

        if not group_option and "LNCEL name" in plot_df.columns:

            for i, cell in enumerate(plot_df["LNCEL name"].unique()):

                cell_df = plot_df[plot_df["LNCEL name"] == cell]

                fig.add_trace(go.Scatter(
                    x=cell_df["Time_str"],
                    y=cell_df[selected_kpi],
                    mode="lines+markers",
                    name=cell,
                    line=dict(color=colors[i % len(colors)]),
                    connectgaps=False
                ))

        else:

            fig.add_trace(go.Scatter(
                x=plot_df["Time_str"],
                y=plot_df[selected_kpi],
                mode="lines+markers",
                name=selected_kpi
            ))

        fig.update_layout(
            height=420,
            title=dict(text=selected_kpi, x=0.5),
            hovermode="x unified",
            legend=dict(x=1.02, y=1),
            margin=dict(l=40, r=120, t=60, b=40)
        )
        if daily_option:
            fig.update_xaxes(tickformat="%Y-%m-%d")
        else:
            fig.update_xaxes(tickformat="%Y-%m-%d %H:%M")

        cols[idx % 2].plotly_chart(fig, use_container_width=True)

        figures.append(fig)


    # -------- PNG EXPORT --------
    export_fig = figures[0]

    # Ensure correct date format
    if daily_option:
        export_fig.update_xaxes(tickformat="%Y-%m-%d")
    else:
        export_fig.update_xaxes(tickformat="%Y-%m-%d %H:%M")

    buf = io.BytesIO()
    export_fig.write_image(buf, format="png", width=1400, height=600, scale=2)

    st.download_button(
        label="📥 Download Chart PNG",
        data=buf.getvalue(),
        file_name="lte_kpi_chart.png",
        mime="image/png"
    )


    # -------- POWERPOINT EXPORT --------
    from pptx.util import Pt  # for font sizes

# -------- POWERPOINT EXPORT --------
def create_ppt(figures):
    # Create a widescreen PPT (16:9)
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16 units
    prs.slide_height = Inches(7.5)   # 9 units

    # Positions for up to 4 charts per slide
    positions = [
        (Inches(0.5), Inches(0.5)),
        (Inches(6.9), Inches(0.5)),
        (Inches(0.5), Inches(4.0)),
        (Inches(6.9), Inches(4.0))
    ]

    # Chart size in inches
    chart_width = Inches(6.08)
    chart_height = Inches(3.04)

    for fig_idx, fig in enumerate(figures):
        # Create a new slide every 4 charts
        if fig_idx % 4 == 0 and fig_idx != 0:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
        elif fig_idx == 0:
            slide = prs.slides.add_slide(prs.slide_layouts[5])

        # Force x-axis to strings for date formatting
        for trace in fig.data:
            if hasattr(trace, 'x'):
                trace.x = [str(x) for x in trace.x]

        # Update layout
        fig.update_layout(
           
            title_font=dict(size=34),   # Increase KPI title size
            margin=dict(l=60, r=150, t=60, b=60),
            legend=dict(
                x=1.05, y=1,
                font=dict(size=28),
                bgcolor="rgba(255,255,255,0)",
                bordercolor="black",
                borderwidth=0.5
            ),
            xaxis=dict(
                title_font=dict(size=32),
                tickfont=dict(size=18) # Smaller for dates
            ),
            yaxis=dict(
                title_font=dict(size=32),
                tickfont=dict(size=32)
            ),
            height=int(chart_height.inches*150),
            width=int(chart_width.inches*300)
        )

        # Export chart as image
        img_buf = io.BytesIO()
        fig.write_image(
            img_buf,
            format="png",
            width=int(chart_width.inches*300),
            height=int(chart_height.inches*300),
            scale=2
        )
        img_buf.seek(0)

        # Determine position on slide
        pos_idx = fig_idx % 4
        slide.shapes.add_picture(
            img_buf,
            positions[pos_idx][0],
            positions[pos_idx][1],
            width=chart_width,
            height=chart_height
        )

    # Save PPT to buffer
    ppt_buffer = io.BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    return ppt_buffer


# ---------- CALL FUNCTION AND DOWNLOAD BUTTON ----------
if figures:
    ppt_file = create_ppt(figures)

    st.download_button(
        label="📊 Download PowerPoint Report",
        data=ppt_file,
        file_name="LTE_KPI_Report.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
else:
    st.warning("⚠️ No data available for the selected filters.")